import base64
from io import BytesIO
from pathlib import Path
import random
import sys
from typing import List
from pptx import Presentation
from pptx.enum.shapes import PP_PLACEHOLDER_TYPE
from pptx.slide import SlideLayout
from dotenv import load_dotenv
load_dotenv()
PROJECT_PATH = Path(__file__).parent.parent.resolve()
sys.path.append(str(PROJECT_PATH))
from src.image_sources import text2image_flux_fp8
from src.image_sources import text2image_sdxl
from src.image_sources import text2image_openai
from src.image_sources import query_image
from src.data_converter import Feature, Image, ImageType, Node, Root, markdown2Root
from src.utils import get_image, valid_path

__all__ = ["generate_ppt_gr"]

def extract_from_layout(slide_layout:SlideLayout):
    placeholder_types: List[int] = []
    placeholder_ids: List[int] = []
    for placeholder in slide_layout.placeholders:
        phf = placeholder.placeholder_format
        if phf.type != PP_PLACEHOLDER_TYPE.DATE and phf.type != PP_PLACEHOLDER_TYPE.FOOTER and phf.type != PP_PLACEHOLDER_TYPE.SLIDE_NUMBER:
            placeholder_types.append(phf.type)
            placeholder_ids.append(phf.idx)
    sorted_ids_types = sorted(zip(placeholder_ids, placeholder_types), key=lambda x: x[1])
    if len(sorted_ids_types) > 0:
        ids, types = zip(*sorted_ids_types)
        return list(ids), list(types)
    return [], []

def types_from_template(template_path:Path):
    prs = Presentation(template_path)
    slide_layouts = prs.slide_layouts
    types = []
    for slide_layout in slide_layouts:
        ids, types = extract_from_layout(slide_layout)
        #print(f"types: {types}")


def features_from_template(template_path:Path) -> dict:
    prs = Presentation(template_path)
    slide_layouts = prs.slide_layouts
    features = {}
    for slide_layout in slide_layouts:
        ids, types = extract_from_layout(slide_layout)
        feature: Feature = Feature()
        feature.set_placeholder_types(types)
        if feature.type not in features:
            #如果不存在则向dict中添加
            features[feature.type] = {"feature":feature, "slide_layout":[slide_layout],"slide_layout_ids":[ids]}
        else:
            #如果存在则向dict中的slide_layout_ids中添加，表示同一种feature存在多种slide_layout
            features[feature.type]["slide_layout"].append(slide_layout)
            features[feature.type]["slide_layout_ids"].append(ids)
    return features

def random_choice(oringin_list:list):
    random_index = random.randint(0, len(oringin_list) - 1)
    random_element = oringin_list[random_index]
    return random_index, random_element

def generate_node(node:Node, features:dict, prs:Presentation, image_source:str):
    idx, slide_layout = random_choice(features[node.feature.type]["slide_layout"])
    feature = features[node.feature.type]["feature"]
    slide = prs.slides.add_slide(slide_layout)
    slide.placeholders[0].text = node.title
    placeholder_types = feature.placeholder_types
    placeholder_ids = features[node.feature.type]["slide_layout_ids"][idx]
    contents = node.contents
    images = node.images
    for index,placeholder_type in enumerate(placeholder_types):
            if placeholder_type == PP_PLACEHOLDER_TYPE.PICTURE:
                image:Image = images.pop(0)
                if image.type == ImageType.KEYWORD:
                    if image_source == "sdxl":
                        img = text2image_sdxl(image.value)
                        buf = BytesIO()
                        img.save(buf, format='PNG')
                        buf.seek(0)
                    elif image_source == "flux":
                        img = text2image_flux_fp8(image.value)
                        buf = BytesIO()
                        img.save(buf, format='PNG')
                        buf.seek(0)
                    elif image_source == "openai":
                        b64_json = text2image_openai(image.value)
                        buf = BytesIO(base64.b64decode(b64_json))
                    elif image_source == "pexels":
                        url = query_image(image.value)
                        buf = BytesIO(get_image(url))
                    else:
                        raise Exception(f"Not support {image_source} image source")
                elif image.type == ImageType.URL:
                    buf = BytesIO(get_image(image.value))
                else:
                    raise Exception(f"Not support {image.type} image type")
                slide.placeholders[placeholder_ids[index]].insert_picture(buf)
            elif placeholder_type == PP_PLACEHOLDER_TYPE.BODY or placeholder_type == PP_PLACEHOLDER_TYPE.OBJECT:
                slide.placeholders[placeholder_ids[index]].text = contents.pop(0)
            elif placeholder_type == PP_PLACEHOLDER_TYPE.SUBTITLE:
                slide.placeholders[placeholder_ids[index]].text = node.subtitle

def generate_children(parent:Node, features:dict, prs:Presentation,image_source:str):
    for node in parent.children:
        generate_node(node, features, prs, image_source)
        generate_children(node, features, prs, image_source)

def generate_ppt(root:Root, template_path:Path,image_source:str) -> Path:
    features = features_from_template(template_path)
    prs = Presentation(template_path)
    #生成主标题
    generate_node(root, features, prs, image_source)
    #生成子标题
    generate_children(root, features, prs, image_source)
    output_path = valid_path(PROJECT_PATH / "tmp/ppt/output.pptx")
    prs.save(output_path)
    return output_path

def generate_ppt_gr(markdown_text:str,theme:str="purple-modern",image_source:str="pexels") -> str:
    template_path = valid_path(PROJECT_PATH / f"src/themes/{theme}.pptx")
    root = markdown2Root(markdown_text)
    return str(generate_ppt(root,template_path,image_source))

